from __future__ import annotations

import os
from pathlib import Path
from typing import List, Dict, Any, Tuple

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pptx import Presentation

from vectorstore import get_text_vectorstore, get_image_vectorstore


def _assets_dir() -> Path:
    p = os.environ.get("CHROMA_ASSETS_DIR")
    if p:
        return Path(p)
    return Path(__file__).resolve().parents[2] / "chroma_assets"


def _load_xlsx_rows(xlsx_path: Path) -> List[Document]:
    """표준개선방안 Excel을 Document로 변환한다(유연 헤더 매핑)."""
    wb = load_workbook(filename=str(xlsx_path), data_only=True)
    docs: List[Document] = []

    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows or len(rows) < 2:
            continue

        header = [str(c).strip() if c is not None else "" for c in rows[0]]

        def idx_of(*cands: str) -> int:
            for c in cands:
                if c in header:
                    return header.index(c)
            return -1

        i_error = idx_of("오류 유형", "오류유형")
        i_check = idx_of("지침")
        i_text = idx_of("문제점 및 개선 방안", "문제점및개선방안")
        i_code = idx_of("참고 사항", "참고사항")

        for r in rows[1:]:
            if not r:
                continue

            def get(i: int) -> str:
                if i < 0 or i >= len(r) or r[i] is None:
                    return ""
                return str(r[i]).strip()

            error_type = get(i_error)
            check_item = get(i_check)
            imp_text = get(i_text)
            imp_code = get(i_code)

            if not (error_type or check_item or imp_text or imp_code):
                continue

            content = "\n".join(
                [
                    f"오류유형: {error_type}",
                    f"검사항목: {check_item}",
                    f"표준개선방안텍스트: {imp_text}",
                    f"표준개선방안코드: {imp_code}",
                ]
            ).strip()

            docs.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": str(xlsx_path),
                        "sheet": ws.title,
                        "kind": "standard",
                        "error_type": error_type,
                        "check_item": check_item,
                    },
                )
            )
    return docs


def _pptx_slide_text(slide) -> str:
    chunks: List[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            txt = shape.text.strip()
            if txt:
                chunks.append(txt)
    return "\n".join(chunks).strip()


def _extract_pptx_images(pptx_path: Path) -> Tuple[List[Document], List[Tuple[str, Dict[str, Any]]]]:
    """PPTX에서 (1) 슬라이드 텍스트 문서, (2) 이미지 파일 경로+메타데이터를 추출"""
    prs = Presentation(str(pptx_path))
    text_docs: List[Document] = []
    image_items: List[Tuple[str, Dict[str, Any]]] = []

    out_dir = _assets_dir() / "pptx_images" / pptx_path.stem
    out_dir.mkdir(parents=True, exist_ok=True)

    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = _pptx_slide_text(slide)
        if slide_text:
            text_docs.append(
                Document(
                    page_content=slide_text,
                    metadata={
                        "source": str(pptx_path),
                        "slide": idx,
                        "kind": "history_text",
                    },
                )
            )

        img_no = 0
        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            try:
                image = shape.image
                blob = image.blob
                ext = image.ext  # 'png', 'jpeg', ...
                img_no += 1
                fname = f"slide_{idx:03d}_img_{img_no:02d}.{ext}"
                fpath = out_dir / fname
                fpath.write_bytes(blob)

                image_items.append(
                    (
                        str(fpath),
                        {
                            "source": str(pptx_path),
                            "slide": idx,
                            "kind": "history_image",
                            # slide_text를 메타로 넣어두면, 검색된 유사 이미지가 어떤 케이스인지 설명하기 쉬움
                            "slide_text": slide_text[:2000] if slide_text else "",
                        },
                    )
                )
            except Exception:
                continue

    return text_docs, image_items


def load_docs(docs_dir: str) -> Tuple[List[Document], List[Tuple[str, Dict[str, Any]]]]:
    p = Path(docs_dir)
    if not p.exists():
        return [], []

    text_docs: List[Document] = []
    image_items: List[Tuple[str, Dict[str, Any]]] = []

    # Excel (표준 개선방안)
    for xlsx in p.rglob("*.xlsx"):
        text_docs.extend(_load_xlsx_rows(xlsx))

    # PPTX (진단 이력: 텍스트 + 이미지)
    for pptx in p.rglob("*.pptx"):
        tdocs, imgs = _extract_pptx_images(pptx)
        text_docs.extend(tdocs)
        image_items.extend(imgs)

    # 기타 txt
    for txt in p.rglob("*.txt"):
        try:
            content = txt.read_text(encoding="utf-8").strip()
            if content:
                text_docs.append(Document(page_content=content, metadata={"source": str(txt), "kind": "text"}))
        except Exception:
            pass

    return text_docs, image_items


def ingest_docs(docs_dir: str = "/app/docs") -> Dict[str, int]:
    """텍스트/이미지를 각각의 Chroma collection에 ingest한다."""
    text_docs, image_items = load_docs(docs_dir)
    out = {"text": 0, "images": 0}

    # 1) Text collection
    if text_docs:
        splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=120)
        split_docs = splitter.split_documents(text_docs)
        tvs = get_text_vectorstore()
        tvs.add_documents(split_docs)
        tvs.persist()
        out["text"] = len(split_docs)

    # 2) Image collection
    if image_items:
        ivs = get_image_vectorstore()
        uris = [u for (u, _m) in image_items]
        metadatas = [_m for (_u, _m) in image_items]
        # add_images가 내부에서 base64로 저장
        ivs.add_images(uris=uris, metadatas=metadatas)
        ivs.persist()
        out["images"] = len(image_items)

    return out


if __name__ == "__main__":
    docs_dir = os.getenv("DOCS_DIR", "/app/docs")
    result = ingest_docs(docs_dir)
    print(f"ingested_text={result['text']} ingested_images={result['images']}")
